# install_and_run.py
# Run: python install_and_run.py

# ── 1. Install dependencies ────────────────────────────────────────────

import os
import json
import requests
from PIL import Image
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()
print("✅ All imports successful!")

# ── 3. PPTGenerator Class ──────────────────────────────────────────────
class PPTGenerator:
    def __init__(self, api_key):
        self.client = Groq(api_key=api_key)
        self.model = "llama-3.3-70b-versatile"
        self.presentation = Presentation()

    def generate_content(self, topic, num_slides=5):
        """Generate content outline using Groq"""
        prompt = f"""
Create a detailed PowerPoint outline on "{topic}" with {num_slides} slides.

Return ONLY valid JSON array in this format:
[
  {{
    "title": "Slide Title",
    "content": "Bullet point 1\\nBullet point 2\\nBullet point 3",
    "slide_type": "title|content|image|conclusion"
  }}
]

Make it engaging and structured.
"""
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "Return ONLY valid JSON. No explanations, no markdown, just pure JSON."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                temperature=0.2,
                response_format={"type": "json_object"}
            )

            content = response.choices[0].message.content.strip()
            parsed = json.loads(content)

            if isinstance(parsed, dict):
                for key in ["slides", "outline", "presentation"]:
                    if key in parsed and isinstance(parsed[key], list):
                        return parsed[key]
                return None

            return parsed

        except Exception as e:
            print(f"❌ Error generating content: {e}")
            return None

    def generate_image_description(self, slide_content):
        """Generate image description based on slide content"""
        prompt = f"""
Based on this slide content, suggest a relevant image description.
Slide content: {slide_content}
Return a concise, descriptive image description (max 30 words).
"""
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {
                        "role": "system",
                        "content": "Return only a short image description, no extra text."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                temperature=0.2
            )
            return response.choices[0].message.content.strip()

        except Exception as e:
            print(f"❌ Error generating image description: {e}")
            return "technology abstract"

    def download_image(self, query, save_path="temp_image.jpg"):
        """Download image from Pexels"""
        try:
            pexels_api_key = "ex67Quw95VyszEN5EwNMhtNIiIa6CpDJvCiWoTyTyJTMLSAI424rjEEf"
            url = "https://api.pexels.com/v1/search"
            headers = {'Authorization': pexels_api_key}
            params = {
                'query': query[:50],
                'per_page': 1,
                'orientation': 'landscape'
            }

            response = requests.get(url, headers=headers, params=params, timeout=10)

            if response.status_code == 200:
                data = response.json()
                if data.get('photos'):
                    image_url = data['photos'][0]['src']['medium']
                    img_response = requests.get(image_url, timeout=10)
                    with open(save_path, 'wb') as f:
                        f.write(img_response.content)
                    return save_path

            # Fallback placeholder image
            img = Image.new('RGB', (800, 600), color='#4A90E2')
            img.save(save_path)
            return save_path

        except Exception as e:
            print(f"❌ Error with image: {e}")
            img = Image.new('RGB', (800, 600), color='#4A90E2')
            img.save(save_path)
            return save_path

    def create_title_slide(self, title, subtitle=""):
        """Create title slide"""
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if subtitle and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = subtitle
            except:
                pass

        return slide

    def create_content_slide(self, title, content, include_image=False):
        """Create content slide with bullet points"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = content

        if include_image:
            try:
                image_desc = self.generate_image_description(content)
                image_path = self.download_image(image_desc)
                if image_path:
                    slide.shapes.add_picture(image_path, Inches(6), Inches(2), height=Inches(4))
            except:
                pass

        return slide

    def generate_presentation(self, topic, num_slides=5, output_path="generated_presentation.pptx"):
        """Generate complete PowerPoint presentation"""
        print(f"📊 Generating presentation on: {topic}")
        print(f"📄 Number of slides: {num_slides}")

        outline = self.generate_content(topic, num_slides)

        if not outline:
            print("❌ Failed to generate content outline")
            return None

        for i, slide_data in enumerate(outline):
            title   = slide_data.get("title", f"Slide {i+1}")
            content = slide_data.get("content", "Content goes here")

            print(f"🎨 Creating slide {i+1}: {title}")

            if i == 0:
                self.create_title_slide(title, "AI Generated Presentation")
            else:
                self.create_content_slide(title, content, include_image=(i % 2 == 0))

        self.presentation.save(output_path)
        print(f"✅ Presentation saved as: {output_path}")
        return output_path

# ── 4. Run ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    GROQ_API_KEY = "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"

    try:
        generator = PPTGenerator(api_key=GROQ_API_KEY)
        print("✅ PPT Generator initialized successfully!")
    except Exception as e:
        print(f"❌ Initialization error: {e}")
        exit()

    topic      = "Artificial Intelligence in Healthcare"
    num_slides = 6
    output     = "ai_healthcare_presentation.pptx"

    output_file = generator.generate_presentation(topic, num_slides, output)

    if output_file:
        print(f"\n🎉 Done! File saved at: {os.path.abspath(output_file)}")